"""pptx 读取 — 按 slide 编号 + 文本框 + 表格 + speaker notes.

只抽文本, 不抽形状几何 / 图片 (那是 file_verify 的活).
"""

from __future__ import annotations

from pathlib import Path

from pentaloom.capabilities.file._models import FileReadResult, FileWarning


def _shape_text(shape) -> str:
    """从一个 shape 取文本. 文本框 / placeholder / 自选图都走这里."""
    if not shape.has_text_frame:
        return ""
    parts: list[str] = []
    for para in shape.text_frame.paragraphs:
        line = "".join(run.text or "" for run in para.runs)
        if line.strip():
            parts.append(line)
    return "\n".join(parts)


def _table_to_markdown(table) -> str:
    rows: list[list[str]] = []
    for row in table.rows:
        rows.append([(cell.text or "").strip().replace("\n", " ") for cell in row.cells])
    if not rows:
        return ""
    width = max(len(r) for r in rows)
    rows = [r + [""] * (width - len(r)) for r in rows]
    lines = ["| " + " | ".join(rows[0]) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for r in rows[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def read_pptx(path: Path) -> FileReadResult:
    from pptx import Presentation  # type: ignore[import-not-found]

    prs = Presentation(str(path))
    sections: list[str] = []
    warnings: list[FileWarning] = []
    empty_slides: list[int] = []

    for i, slide in enumerate(prs.slides, start=1):
        body_parts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                t = _shape_text(shape)
                if t:
                    body_parts.append(t)
            if getattr(shape, "has_table", False):
                md = _table_to_markdown(shape.table)
                if md:
                    body_parts.append(md)

        notes = ""
        if slide.has_notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf is not None:
                notes = (notes_tf.text or "").strip()

        if not body_parts and not notes:
            empty_slides.append(i)
            sections.append(f"## Slide {i}\n_(empty)_")
            continue

        chunk = [f"## Slide {i}"]
        if body_parts:
            chunk.append("\n".join(body_parts))
        if notes:
            chunk.append(f"\n_Notes:_\n{notes}")
        sections.append("\n".join(chunk))

    text = "\n\n".join(sections).strip()
    meta: dict[str, int | str | list[str]] = {
        "slides": len(prs.slides),
        "chars": len(text),
    }
    if empty_slides:
        meta["empty_slides"] = empty_slides
        warnings.append(FileWarning(
            kind="empty_slide",
            message=f"{len(empty_slides)} 张空 slide (编号: {empty_slides})",
        ))

    return FileReadResult(
        path=str(path),
        kind="pptx",
        text=text,
        warnings=warnings,
        meta=meta,
    )
